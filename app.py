import streamlit as st
import pandas as pd
import numpy as np
from io import StringIO, BytesIO
import matplotlib.pyplot as plt
import seaborn as sns
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches
from docx import Document
from docx.shared import Inches as DocxInches
import google.generativeai as genai
import os
import tempfile

# ---------------------------------------------------
# CONFIGURE GEMINI API KEY
# ---------------------------------------------------
GEMINI_API_KEY = "AIzaSyBGGZMnCKXGNCkdO3UpElZ-81WibZozsuM"  # 👈 Your Gemini key
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel('gemini-1.5-flash')

# ---------------------------------------------------
# PAGE CONFIGURATION
# ---------------------------------------------------
st.set_page_config(
    page_title="Nexus AI Analytics",
    page_icon="📊",
    layout="wide"
)

# ---------------------------------------------------
# PROFESSIONAL LIGHT UI
# ---------------------------------------------------
st.markdown("""
<style>
.stApp { background-color: #F3F4F6; }
section[data-testid="stSidebar"] { background-color: #E5E7EB; }
html, body, [class*="css"] { color: #111827; font-family: 'Segoe UI', sans-serif; }
h1, h2, h3, h4, h5 { color: #111827 !important; }
p { color: #374151 !important; }
label { color: #111827 !important; font-weight: 600; }
.card {
    background: white; padding: 25px; border-radius: 18px;
    border: 1px solid #E5E7EB; box-shadow: 0 4px 12px rgba(0,0,0,0.08);
    margin-bottom: 20px;
}
[data-testid="metric-container"] {
    background: white; border: 1px solid #E5E7EB;
    padding: 18px; border-radius: 16px;
    box-shadow: 0 2px 10px rgba(0,0,0,0.05);
}
[data-testid="stFileUploader"] {
    background: white; border: 2px dashed #CBD5E1;
    border-radius: 18px; padding: 20px;
}
button[data-baseweb="tab"] {
    background: #E5E7EB; color: #111827 !important;
    border-radius: 10px; margin-right: 5px;
    padding: 10px 18px; font-weight: 600;
}
.stButton button {
    background: linear-gradient(to right, #2563EB, #1D4ED8);
    color: white !important; border: none;
    border-radius: 12px; padding: 12px 24px; font-weight: 600;
}
.stDownloadButton button {
    background: linear-gradient(to right, #10B981, #059669);
    color: white !important; border: none;
    border-radius: 12px; padding: 12px 24px; font-weight: 600;
}
input, textarea {
    background-color: white !important;
    color: #111827 !important; border-radius: 10px !important;
}
div[data-baseweb="select"] { background-color: white; border-radius: 10px; }
[data-testid="stDataFrame"] { background-color: white; }
[data-testid="stAlert"] { border-radius: 12px; }
</style>
""", unsafe_allow_html=True)

# ---------------------------------------------------
# SESSION STATE
# ---------------------------------------------------
if "df" not in st.session_state:
    st.session_state.df = None
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "chart_paths" not in st.session_state:
    st.session_state.chart_paths = []

# ---------------------------------------------------
# HELPER FUNCTIONS
# ---------------------------------------------------
def clean_dataframe(df):
    actions = []
    cleaned = df.copy()
    before = len(cleaned)
    cleaned = cleaned.drop_duplicates()
    actions.append(f"✅ Removed {before - len(cleaned)} duplicate rows")
    
    for col in cleaned.columns:
        missing = cleaned[col].isnull().sum()
        if missing > 0:
            if cleaned[col].dtype in [np.float64, np.int64]:
                cleaned[col].fillna(cleaned[col].median(), inplace=True)
                actions.append(f"✅ Filled {missing} missing in '{col}' with median")
            else:
                mode_val = cleaned[col].mode()[0] if not cleaned[col].mode().empty else "Unknown"
                cleaned[col].fillna(mode_val, inplace=True)
                actions.append(f"✅ Filled {missing} missing in '{col}' with mode")
    
    for col in cleaned.select_dtypes(include="object").columns:
        cleaned[col] = cleaned[col].astype(str).str.strip()
    actions.append("✅ Trimmed whitespace from text columns")
    return cleaned, actions


def generate_charts(df):
    chart_paths = []
    temp_dir = tempfile.gettempdir()
    numeric_cols = df.select_dtypes(include=np.number).columns.tolist()
    cat_cols = df.select_dtypes(include="object").columns.tolist()
    sns.set_style("whitegrid")
    
    for col in numeric_cols[:3]:
        plt.figure(figsize=(8, 5))
        sns.histplot(df[col].dropna(), kde=True, color="#2563EB")
        plt.title(f"Distribution of {col}", fontsize=14, fontweight='bold')
        path = os.path.join(temp_dir, f"hist_{col}.png")
        plt.savefig(path, bbox_inches="tight", dpi=100)
        plt.close()
        chart_paths.append((f"Distribution of {col}", path))
    
    if len(numeric_cols) > 1:
        plt.figure(figsize=(10, 6))
        sns.heatmap(df[numeric_cols].corr(), annot=True, cmap="coolwarm", fmt=".2f")
        plt.title("Correlation Heatmap", fontsize=14, fontweight='bold')
        path = os.path.join(temp_dir, "heatmap.png")
        plt.savefig(path, bbox_inches="tight", dpi=100)
        plt.close()
        chart_paths.append(("Correlation Heatmap", path))
    
    for col in cat_cols[:2]:
        if df[col].nunique() < 15:
            plt.figure(figsize=(8, 5))
            df[col].value_counts().head(10).plot(kind="bar", color="#10B981")
            plt.title(f"Count of {col}", fontsize=14, fontweight='bold')
            plt.xticks(rotation=45)
            path = os.path.join(temp_dir, f"bar_{col}.png")
            plt.savefig(path, bbox_inches="tight", dpi=100)
            plt.close()
            chart_paths.append((f"Count of {col}", path))
    return chart_paths


def generate_word_report(df, chart_paths, ai_insights=""):
    doc = Document()
    doc.add_heading("Nexus AI Analytics Report", 0)
    doc.add_paragraph("Generated by Nexus AI Analytics Platform")
    
    doc.add_heading("Dataset Overview", level=1)
    doc.add_paragraph(f"Total Rows: {len(df)}")
    doc.add_paragraph(f"Total Columns: {len(df.columns)}")
    doc.add_paragraph(f"Numeric Columns: {len(df.select_dtypes(include=np.number).columns)}")
    doc.add_paragraph(f"Missing Values: {df.isnull().sum().sum()}")
    
    doc.add_heading("Summary Statistics", level=1)
    doc.add_paragraph(df.describe().to_string())
    
    doc.add_heading("Column Information", level=1)
    for col in df.columns:
        doc.add_paragraph(f"• {col} ({df[col].dtype}) - Missing: {df[col].isnull().sum()}", style='List Bullet')
    
    if chart_paths:
        doc.add_heading("Visualizations", level=1)
        for title, path in chart_paths:
            doc.add_heading(title, level=2)
            doc.add_picture(path, width=DocxInches(6))
    
    if ai_insights:
        doc.add_heading("AI Insights", level=1)
        doc.add_paragraph(ai_insights)
    
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer


def generate_ppt(df, chart_paths, ai_insights=""):
    prs = Presentation()
    
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Nexus AI Analytics Report"
    slide.placeholders[1].text = "Auto-generated by AI Data Analyst"
    
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Dataset Overview"
    tf = slide.placeholders[1].text_frame
    tf.text = f"Rows: {len(df)}"
    tf.add_paragraph().text = f"Columns: {len(df.columns)}"
    tf.add_paragraph().text = f"Numeric Columns: {len(df.select_dtypes(include=np.number).columns)}"
    tf.add_paragraph().text = f"Missing Values: {df.isnull().sum().sum()}"
    
    for title, path in chart_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_picture(path, Inches(1), Inches(1.5), Inches(8), Inches(5))
    
    if ai_insights:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "AI Insights"
        slide.placeholders[1].text = ai_insights[:1500]
    
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def get_ai_insights(df):
    try:
        summary = f"""
        Dataset Shape: {df.shape}
        Columns: {list(df.columns)}
        Data Types: {df.dtypes.to_dict()}
        Statistics: {df.describe().to_string()}
        Sample: {df.head(3).to_dict()}
        """
        prompt = f"You are a senior data analyst. Analyze this dataset and provide 5-7 key business insights with actionable recommendations:\n{summary}"
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"AI insights unavailable. Error: {str(e)}\n\nBasic insights:\n• Dataset has {len(df)} rows and {len(df.columns)} columns\n• Missing values: {df.isnull().sum().sum()}"


def chat_with_ai(message, df):
    try:
        context = ""
        if df is not None:
            context = f"""
            User's data:
            - Shape: {df.shape}
            - Columns: {list(df.columns)}
            - Sample: {df.head(3).to_dict()}
            - Statistics: {df.describe().to_string()}
            """
        prompt = f"You are an AI Data Analyst assistant. {context}\n\nUser question: {message}"
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        return f"Error: {str(e)}"


# ---------------------------------------------------
# SIDEBAR
# ---------------------------------------------------
st.sidebar.title("🚀 Nexus AI Analytics")
st.sidebar.markdown("---")
st.sidebar.subheader("📌 How To Use")
st.sidebar.info("""
1. Select a data source
2. Upload or paste dataset
3. Click Run AI Analysis
4. Clean data automatically
5. Generate charts & reports
6. Chat with AI assistant
""")
st.sidebar.markdown("---")

# AI Chatbot
st.sidebar.subheader("💬 AI Assistant")
user_input = st.sidebar.text_input("Ask about your data:", key="chat_input")
if st.sidebar.button("Send 💬"):
    if user_input:
        st.session_state.chat_history.append(("user", user_input))
        with st.spinner("Thinking..."):
            reply = chat_with_ai(user_input, st.session_state.df)
        st.session_state.chat_history.append(("bot", reply))

if st.session_state.chat_history:
    st.sidebar.markdown("### Conversation")
    for role, msg in st.session_state.chat_history[-6:]:
        if role == "user":
            st.sidebar.markdown(f"**🧑 You:** {msg}")
        else:
            st.sidebar.markdown(f"**🤖 AI:** {msg}")
    if st.sidebar.button("Clear Chat"):
        st.session_state.chat_history = []
        st.rerun()

st.sidebar.markdown("---")
st.sidebar.success("System Ready ✅")

# ---------------------------------------------------
# MAIN
# ---------------------------------------------------
st.title("📊 Nexus AI Analytics")
st.markdown("Professional AI-powered business intelligence dashboard. Powered by Google Gemini.")
st.markdown("---")

# FEATURE CARDS
col1, col2, col3, col4 = st.columns(4)
with col1:
    st.markdown('<div class="card"><h3>📂 Smart Input</h3><p>CSV/Excel upload, paste, or URLs.</p></div>', unsafe_allow_html=True)
with col2:
    st.markdown('<div class="card"><h3>🧹 Auto Cleaning</h3><p>Remove duplicates & fill missing values.</p></div>', unsafe_allow_html=True)
with col3:
    st.markdown('<div class="card"><h3>📊 Visual Analytics</h3><p>Charts, heatmaps, and visualizations.</p></div>', unsafe_allow_html=True)
with col4:
    st.markdown('<div class="card"><h3>🤖 AI Reports</h3><p>Word & PPT auto-generation.</p></div>', unsafe_allow_html=True)

st.markdown("---")

# DATA INPUT
st.subheader("📂 Smart Data Input Center")
input_method = st.selectbox(
    "Choose Data Source",
    ["📁 Upload CSV/Excel File", "📋 Paste Sample Dataset", "🔗 Load Dataset From URL"]
)

if input_method == "📁 Upload CSV/Excel File":
    uploaded_file = st.file_uploader("Upload Dataset", type=["csv", "xlsx", "xls"])
    if uploaded_file is not None:
        try:
            if uploaded_file.name.endswith('.csv'):
                st.session_state.df = pd.read_csv(uploaded_file)
            else:
                st.session_state.df = pd.read_excel(uploaded_file)
            st.success("✅ Dataset uploaded successfully")
        except Exception as e:
            st.error(f"Error: {e}")

elif input_method == "📋 Paste Sample Dataset":
    sample_data = st.text_area(
        "Paste CSV Data Below", height=250,
        placeholder="Order_ID,Region,Sales,Profit\n1001,North,5000,1200\n1002,South,7000,1800"
    )
    if sample_data:
        try:
            st.session_state.df = pd.read_csv(StringIO(sample_data))
            st.success("✅ Sample dataset loaded successfully")
        except Exception as e:
            st.error(f"Error: {e}")

elif input_method == "🔗 Load Dataset From URL":
    file_url = st.text_input("Paste CSV File URL")
    if file_url:
        try:
            st.session_state.df = pd.read_csv(file_url)
            st.success("✅ Dataset loaded successfully")
        except Exception as e:
            st.error(f"Unable to load: {e}")

# BUTTONS
col1, col2, col3 = st.columns(3)
with col1:
    run_analysis = st.button("🚀 Run AI Analysis")
with col2:
    clean_btn = st.button("🧹 Auto Clean Data")
with col3:
    generate_reports = st.button("📄 Generate Reports")

# CLEAN
if clean_btn and st.session_state.df is not None:
    with st.spinner("Cleaning data..."):
        cleaned, actions = clean_dataframe(st.session_state.df)
        st.session_state.df = cleaned
    st.markdown("### 🧹 Data Cleaning Report")
    for action in actions:
        st.write(action)
    st.success("✅ Data cleaned successfully!")

# ANALYSIS
if run_analysis and st.session_state.df is not None:
    df = st.session_state.df
    st.markdown("---")
    st.subheader("📌 Dashboard Overview")
    
    col1, col2, col3, col4 = st.columns(4)
    col1.metric("Rows", len(df))
    col2.metric("Columns", len(df.columns))
    col3.metric("Missing Values", df.isnull().sum().sum())
    col4.metric("Numeric Columns", len(df.select_dtypes(include='number').columns))
    
    st.markdown("---")
    
    tab1, tab2, tab3, tab4 = st.tabs(["📄 Dataset", "📊 Visual Analytics", "🤖 AI Insights", "📈 Advanced Charts"])
    
    with tab1:
        st.subheader("Dataset Preview")
        st.dataframe(df, use_container_width=True)
        st.subheader("Statistical Summary")
        st.dataframe(df.describe(), use_container_width=True)
        st.subheader("Column Information")
        info_df = pd.DataFrame({
            'Column': df.columns,
            'Type': df.dtypes.astype(str),
            'Missing': df.isnull().sum().values,
            'Unique': [df[col].nunique() for col in df.columns]
        })
        st.dataframe(info_df, use_container_width=True)
    
    with tab2:
        st.subheader("Visual Analytics")
        numeric_df = df.select_dtypes(include='number')
        if not numeric_df.empty:
            st.markdown("### 📈 Line Chart")
            st.line_chart(numeric_df)
            st.markdown("### 📊 Bar Chart")
            st.bar_chart(numeric_df)
            st.markdown("### 📉 Area Chart")
            st.area_chart(numeric_df)
        else:
            st.warning("No numeric columns available.")
    
    with tab3:
        st.subheader("🤖 AI Generated Insights")
        with st.spinner("Generating AI insights with Gemini..."):
            ai_insights = get_ai_insights(df)
        st.markdown(ai_insights)
        st.download_button(
            label="⬇ Download AI Insights",
            data=ai_insights,
            file_name="AI_Insights.txt",
            mime="text/plain"
        )
    
    with tab4:
        st.subheader("📈 Advanced Visualizations")
        with st.spinner("Generating charts..."):
            chart_paths = generate_charts(df)
            st.session_state.chart_paths = chart_paths
        for title, path in chart_paths:
            st.markdown(f"### {title}")
            st.image(path)
        
        numeric_cols = df.select_dtypes(include=np.number).columns.tolist()
        if len(numeric_cols) >= 2:
            st.markdown("### 🎯 Interactive Scatter Plot")
            x_col = st.selectbox("X-axis", numeric_cols, key="x_axis")
            y_col = st.selectbox("Y-axis", numeric_cols, index=1, key="y_axis")
            fig = px.scatter(df, x=x_col, y=y_col, title=f"{x_col} vs {y_col}")
            st.plotly_chart(fig, use_container_width=True)

# REPORTS
if generate_reports and st.session_state.df is not None:
    df = st.session_state.df
    st.markdown("---")
    st.subheader("📄 Download Reports")
    
    with st.spinner("Generating reports..."):
        if not st.session_state.chart_paths:
            st.session_state.chart_paths = generate_charts(df)
        ai_insights = get_ai_insights(df)
        word_buffer = generate_word_report(df, st.session_state.chart_paths, ai_insights)
        ppt_buffer = generate_ppt(df, st.session_state.chart_paths, ai_insights)
    
    col1, col2, col3 = st.columns(3)
    with col1:
        st.download_button(
            label="📄 Download Word Report",
            data=word_buffer,
            file_name="Nexus_AI_Report.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
    with col2:
        st.download_button(
            label="🎯 Download PowerPoint",
            data=ppt_buffer,
            file_name="Nexus_AI_Report.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    with col3:
        csv_buffer = df.to_csv(index=False).encode('utf-8')
        st.download_button(
            label="📊 Download Cleaned CSV",
            data=csv_buffer,
            file_name="Cleaned_Dataset.csv",
            mime="text/csv"
        )
    st.success("✅ All reports generated successfully!")

st.markdown("---")
st.caption("Nexus AI Analytics Platform | Powered by Streamlit & Google Gemini")
